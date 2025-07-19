"""
PowerPoint metadata extractor using python-pptx library.
Extracts comprehensive metadata about all objects, their formatting properties, and positions.
"""

import os
import json
from typing import Dict, List, Any, Optional, Tuple, Union
from datetime import datetime
from pathlib import Path
import logging

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.picture import Picture
    from pptx.table import Table
    from pptx.shapes.graphfrm import GraphicFrame
    from pptx.shapes.connector import Connector
    from pptx.shapes.freeform import FreeformBuilder
    from pptx.shapes.group import GroupShape
    from pptx.text.text import TextFrame
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
    print("python-pptx imported successfully")
except ImportError as e:
    PPTX_AVAILABLE = False
    print(f"python-pptx import failed: {e}")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PowerPointMetadataExtractor:
    """
    Comprehensive PowerPoint metadata extractor using python-pptx.
    Extracts all objects, their formatting properties, and positions on slides.
    """
    
    def __init__(self, presentation_path: Optional[str] = None):
        """
        Initialize the PowerPoint metadata extractor.
        
        Args:
            presentation_path: Path to the PowerPoint file (optional, can be set later)
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required but not installed. Please install with: pip install python-pptx")
            
        self.presentation_path = presentation_path
        self.presentation = None
        
    def open_presentation(self, presentation_path: Optional[str] = None) -> None:
        """
        Open the PowerPoint presentation.
        
        Args:
            presentation_path: Path to the PowerPoint file
        """
        if presentation_path:
            self.presentation_path = presentation_path
            
        if not self.presentation_path:
            raise ValueError("No presentation path specified")
            
        if not os.path.exists(self.presentation_path):
            raise FileNotFoundError(f"Presentation not found: {self.presentation_path}")
            
        try:
            self.presentation = Presentation(self.presentation_path)
        except Exception as e:
            raise Exception(f"Failed to open presentation: {str(e)}")
    
    def close(self) -> None:
        """Close the presentation and clean up resources."""
        self.presentation = None
    
    def __enter__(self):
        """Context manager entry."""
        return self
        
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit - cleanup resources."""
        self.close()
        
    def extract_presentation_metadata(
        self, 
        presentation_path: Optional[str] = None,
        include_slide_content: bool = True,
        include_master_slides: bool = True,
        include_layouts: bool = True
    ) -> Dict[str, Any]:
        """
        Extract comprehensive metadata from the PowerPoint presentation.
        
        Args:
            presentation_path: Path to the PowerPoint file
            include_slide_content: Whether to include detailed slide content
            include_master_slides: Whether to include slide master information
            include_layouts: Whether to include slide layout information
            
        Returns:
            Dictionary containing complete presentation metadata
        """
        if presentation_path or not self.presentation:
            self.open_presentation(presentation_path)
            
        try:
            # Get basic presentation info
            presentation_metadata = {
                "extractedAt": datetime.now().isoformat(),
                "presentationPath": self.presentation_path,
                "presentationName": os.path.basename(self.presentation_path) if self.presentation_path else "Unknown",
                "coreProperties": self._extract_core_properties(),
                "slideSize": self._extract_slide_size(),
                "totalSlides": len(self.presentation.slides),
                "slideMasters": [],
                "slideLayouts": [],
                "slides": []
            }
            
            # Extract slide masters if requested
            if include_master_slides:
                presentation_metadata["slideMasters"] = self._extract_slide_masters()
                
            # Extract slide layouts if requested  
            if include_layouts:
                presentation_metadata["slideLayouts"] = self._extract_slide_layouts()
            
            # Extract slide content if requested
            if include_slide_content:
                for i, slide in enumerate(self.presentation.slides):
                    slide_metadata = self._extract_slide_metadata(slide, i + 1)
                    presentation_metadata["slides"].append(slide_metadata)
                    
            return presentation_metadata
            
        except Exception as e:
            raise Exception(f"Failed to extract presentation metadata: {str(e)}")
    
    def _extract_core_properties(self) -> Dict[str, Any]:
        """Extract core presentation properties."""
        try:
            core_props = self.presentation.core_properties
            return {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "comments": core_props.comments or "",
                "category": core_props.category or "",
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "lastModifiedBy": core_props.last_modified_by or "",
                "revision": core_props.revision,
                "version": core_props.version or ""
            }
        except Exception as e:
            logger.warning(f"Error extracting core properties: {str(e)}")
            return {"error": str(e)}
    
    def _extract_slide_size(self) -> Dict[str, Any]:
        """Extract slide size information."""
        try:
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            return {
                "width": slide_width,
                "height": slide_height,
                "widthInches": slide_width / 914400,  # EMU to inches
                "heightInches": slide_height / 914400,
                "aspectRatio": slide_width / slide_height if slide_height > 0 else 0
            }
        except Exception as e:
            logger.warning(f"Error extracting slide size: {str(e)}")
            return {"error": str(e)}
    
    def _extract_slide_masters(self) -> List[Dict[str, Any]]:
        """Extract slide master information."""
        masters = []
        try:
            for i, master in enumerate(self.presentation.slide_masters):
                master_data = {
                    "masterIndex": i,
                    "name": getattr(master, 'name', f"Master_{i}"),
                    "shapes": self._extract_shapes(master.shapes)
                }
                masters.append(master_data)
        except Exception as e:
            logger.warning(f"Error extracting slide masters: {str(e)}")
            masters.append({"error": str(e)})
        return masters
    
    def _extract_slide_layouts(self) -> List[Dict[str, Any]]:
        """Extract slide layout information."""
        layouts = []
        try:
            for master in self.presentation.slide_masters:
                for i, layout in enumerate(master.slide_layouts):
                    layout_data = {
                        "layoutIndex": i,
                        "name": getattr(layout, 'name', f"Layout_{i}"),
                        "shapes": self._extract_shapes(layout.shapes)
                    }
                    layouts.append(layout_data)
        except Exception as e:
            logger.warning(f"Error extracting slide layouts: {str(e)}")
            layouts.append({"error": str(e)})
        return layouts
    
    def _extract_slide_metadata(self, slide, slide_number: int) -> Dict[str, Any]:
        """
        Extract comprehensive metadata from a single slide.
        
        Args:
            slide: python-pptx Slide object
            slide_number: 1-based slide number
            
        Returns:
            Dictionary containing slide metadata
        """
        try:
            slide_metadata = {
                "slideNumber": slide_number,
                "slideId": slide.slide_id,
                "name": getattr(slide, 'name', f"Slide_{slide_number}"),
                "layoutName": getattr(slide.slide_layout, 'name', "Unknown Layout"),
                "shapes": self._extract_shapes(slide.shapes),
                "notes": self._extract_slide_notes(slide),
                "comments": self._extract_slide_comments(slide)
            }
            return slide_metadata
            
        except Exception as e:
            logger.warning(f"Error extracting metadata for slide {slide_number}: {str(e)}")
            return {
                "slideNumber": slide_number,
                "error": str(e)
            }
    
    def _extract_shapes(self, shapes) -> List[Dict[str, Any]]:
        """Extract comprehensive information about all shapes on a slide."""
        shape_list = []
        
        for i, shape in enumerate(shapes):
            try:
                shape_data = self._extract_single_shape(shape, i)
                shape_list.append(shape_data)
            except Exception as e:
                logger.warning(f"Error extracting shape {i}: {str(e)}")
                shape_list.append({
                    "shapeIndex": i,
                    "error": str(e)
                })
                
        return shape_list
    
    def _extract_single_shape(self, shape, shape_index: int) -> Dict[str, Any]:
        """Extract comprehensive metadata from a single shape."""
        base_data = {
            "shapeIndex": shape_index,
            "shapeId": shape.shape_id,
            "name": shape.name,
            "shapeType": str(shape.shape_type),
            "shapeTypeValue": int(shape.shape_type),
            "position": self._extract_position(shape),
            "rotation": getattr(shape, 'rotation', 0),
            "visible": getattr(shape, 'visible', True),
            "zOrder": getattr(shape, 'z_order', 0)
        }
        
        # Add shape-specific data based on type
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
            base_data["textContent"] = self._extract_text_content(shape)
            
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            base_data["imageData"] = self._extract_image_data(shape)
            
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            base_data["tableData"] = self._extract_table_data(shape)
            
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            base_data["autoShapeData"] = self._extract_autoshape_data(shape)
            
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            base_data["groupData"] = self._extract_group_data(shape)
            
        # Extract common formatting
        base_data["fill"] = self._extract_fill_format(shape)
        base_data["line"] = self._extract_line_format(shape) 
        base_data["shadow"] = self._extract_shadow_format(shape)
        
        return base_data
    
    def _extract_position(self, shape) -> Dict[str, Any]:
        """Extract position and size information from a shape."""
        try:
            return {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "leftInches": shape.left / 914400,  # EMU to inches
                "topInches": shape.top / 914400,
                "widthInches": shape.width / 914400,
                "heightInches": shape.height / 914400
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_text_content(self, shape) -> Dict[str, Any]:
        """Extract comprehensive text content and formatting."""
        if not shape.has_text_frame:
            return {"hasText": False}
            
        try:
            text_frame = shape.text_frame
            text_data = {
                "hasText": True,
                "text": text_frame.text,
                "autoSize": str(text_frame.auto_size),
                "margins": {
                    "left": text_frame.margin_left,
                    "right": text_frame.margin_right,
                    "top": text_frame.margin_top,
                    "bottom": text_frame.margin_bottom
                },
                "wordWrap": text_frame.word_wrap,
                "paragraphs": []
            }
            
            # Extract paragraph-level formatting
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                para_data = {
                    "paragraphIndex": para_idx,
                    "text": paragraph.text,
                    "level": paragraph.level,
                    "alignment": str(paragraph.alignment),
                    "lineSpacing": getattr(paragraph, 'line_spacing', None),
                    "spaceAfter": getattr(paragraph, 'space_after', None),
                    "spaceBefore": getattr(paragraph, 'space_before', None),
                    "runs": []
                }
                
                # Extract run-level formatting
                for run_idx, run in enumerate(paragraph.runs):
                    run_data = {
                        "runIndex": run_idx,
                        "text": run.text,
                        "font": self._extract_font_format(run.font)
                    }
                    para_data["runs"].append(run_data)
                
                text_data["paragraphs"].append(para_data)
                
            return text_data
            
        except Exception as e:
            return {"hasText": True, "error": str(e)}
    
    def _extract_font_format(self, font) -> Dict[str, Any]:
        """Extract font formatting information."""
        try:
            font_data = {
                "languageId": getattr(font, 'language_id', None)
            }
            
            # Safely extract each property
            try:
                font_data["name"] = font.name
            except:
                font_data["name"] = None
                
            try:
                font_data["size"] = font.size.pt if font.size else None
            except:
                font_data["size"] = None
                
            try:
                font_data["bold"] = font.bold
            except:
                font_data["bold"] = None
                
            try:
                font_data["italic"] = font.italic
            except:
                font_data["italic"] = None
                
            try:
                font_data["underline"] = font.underline
            except:
                font_data["underline"] = None
                
            try:
                font_data["color"] = self._extract_color(font.color)
            except Exception as color_error:
                font_data["color"] = {"error": str(color_error)}
                
            return font_data
            
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_color(self, color_format) -> Dict[str, Any]:
        """Extract color information."""
        try:
            if not color_format or not hasattr(color_format, 'type'):
                return {"type": "none"}
                
            color_data = {
                "type": str(color_format.type),
                "typeValue": int(color_format.type)
            }
            
            # Handle RGB colors
            if hasattr(color_format, 'rgb') and color_format.type == MSO_COLOR_TYPE.RGB:
                try:
                    color_data["rgb"] = str(color_format.rgb)
                except:
                    color_data["rgb"] = "unknown"
            
            # Handle theme colors - check if THEME constant exists
            elif hasattr(MSO_COLOR_TYPE, 'THEME') and color_format.type == MSO_COLOR_TYPE.THEME:
                try:
                    color_data["themeColor"] = str(color_format.theme_color)
                except:
                    color_data["themeColor"] = "unknown"
            
            # Handle scheme colors
            elif hasattr(color_format, 'theme_color'):
                try:
                    color_data["themeColor"] = str(color_format.theme_color)
                except:
                    color_data["themeColor"] = "unknown"
                
            return color_data
            
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_image_data(self, shape) -> Dict[str, Any]:
        """Extract image-specific data."""
        try:
            if not isinstance(shape, Picture):
                return {"error": "Not a picture shape"}
                
            image = shape.image
            return {
                "filename": image.filename,
                "contentType": image.content_type,
                "size": len(image.blob),
                "dimensions": {
                    "width": getattr(image, 'width', None),
                    "height": getattr(image, 'height', None)
                },
                "crop": {
                    "left": getattr(shape, 'crop_left', 0),
                    "top": getattr(shape, 'crop_top', 0),
                    "right": getattr(shape, 'crop_right', 0),
                    "bottom": getattr(shape, 'crop_bottom', 0)
                }
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_table_data(self, shape) -> Dict[str, Any]:
        """Extract table-specific data."""
        try:
            if not isinstance(shape, Table):
                return {"error": "Not a table shape"}
                
            table = shape.table
            table_data = {
                "rows": len(table.rows),
                "columns": len(table.columns),
                "cells": []
            }
            
            # Extract cell data
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_data = {
                        "row": row_idx,
                        "column": col_idx,
                        "text": cell.text,
                        "textFrame": self._extract_text_content(cell) if cell.text_frame else None,
                        "fill": self._extract_fill_format(cell),
                        "margins": {
                            "left": cell.margin_left,
                            "right": cell.margin_right, 
                            "top": cell.margin_top,
                            "bottom": cell.margin_bottom
                        }
                    }
                    table_data["cells"].append(cell_data)
                    
            return table_data
            
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_autoshape_data(self, shape) -> Dict[str, Any]:
        """Extract autoshape-specific data."""
        try:
            return {
                "autoShapeType": str(getattr(shape, 'auto_shape_type', 'Unknown')),
                "adjustments": getattr(shape, 'adjustments', [])
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_group_data(self, shape) -> Dict[str, Any]:
        """Extract group shape data."""
        try:
            if not isinstance(shape, GroupShape):
                return {"error": "Not a group shape"}
                
            return {
                "shapeCount": len(shape.shapes),
                "shapes": self._extract_shapes(shape.shapes)
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_fill_format(self, shape_or_cell) -> Dict[str, Any]:
        """Extract fill formatting information."""
        try:
            if not hasattr(shape_or_cell, 'fill'):
                return {"type": "none"}
                
            fill = shape_or_cell.fill
            fill_data = {
                "type": str(fill.type) if hasattr(fill, 'type') else "unknown"
            }
            
            # Add type-specific data with error handling
            try:
                if hasattr(fill, 'fore_color'):
                    fill_data["foreColor"] = self._extract_color(fill.fore_color)
            except Exception as color_error:
                fill_data["foreColor"] = {"error": str(color_error)}
                
            try:
                if hasattr(fill, 'back_color'):
                    fill_data["backColor"] = self._extract_color(fill.back_color)
            except Exception as color_error:
                fill_data["backColor"] = {"error": str(color_error)}
                
            return fill_data
            
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_line_format(self, shape) -> Dict[str, Any]:
        """Extract line/border formatting information."""
        try:
            if not hasattr(shape, 'line'):
                return {"type": "none"}
                
            line = shape.line
            return {
                "color": self._extract_color(line.color) if hasattr(line, 'color') else None,
                "width": getattr(line, 'width', None),
                "dashStyle": str(getattr(line, 'dash_style', 'solid'))
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_shadow_format(self, shape) -> Dict[str, Any]:
        """Extract shadow formatting information."""
        try:
            if not hasattr(shape, 'shadow'):
                return {"visible": False}
                
            shadow = shape.shadow
            return {
                "visible": getattr(shadow, 'visible', False),
                "style": str(getattr(shadow, 'style', 'none'))
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_slide_notes(self, slide) -> Dict[str, Any]:
        """Extract slide notes information."""
        try:
            if not hasattr(slide, 'notes_slide') or not slide.notes_slide:
                return {"hasNotes": False}
                
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                return {
                    "hasNotes": True,
                    "text": notes_slide.notes_text_frame.text,
                    "textContent": self._extract_text_content(notes_slide.notes_text_frame)
                }
            else:
                return {"hasNotes": False}
                
        except Exception as e:
            return {"hasNotes": False, "error": str(e)}
    
    def _extract_slide_comments(self, slide) -> List[Dict[str, Any]]:
        """Extract slide comments information."""
        try:
            # Note: python-pptx doesn't directly support comments extraction
            # This is a placeholder for future implementation
            return []
        except Exception as e:
            return [{"error": str(e)}]
    
    def extract_metadata(self, presentation_path: str) -> Dict[str, Any]:
        """
        Extract metadata for cache handler compatibility.
        
        Args:
            presentation_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing presentation metadata
        """
        return self.extract_presentation_metadata(presentation_path)
    
    def extract_to_json(
        self,
        presentation_path: Optional[str] = None,
        output_path: Optional[str] = None,
        **kwargs
    ) -> str:
        """
        Extract metadata and return as JSON string.
        
        Args:
            presentation_path: Path to the PowerPoint file
            output_path: Optional path to save JSON file
            **kwargs: Additional arguments for extract_presentation_metadata
            
        Returns:
            JSON string containing the metadata
        """
        metadata = self.extract_presentation_metadata(presentation_path, **kwargs)
        json_str = json.dumps(metadata, indent=2, ensure_ascii=False)
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(json_str)
            logger.info(f"Metadata saved to: {output_path}")
            
        return json_str


def main():
    """Example usage of the PowerPoint metadata extractor."""
    # This is for testing/demonstration purposes
    extractor = PowerPointMetadataExtractor()
    
    # Example of what can be extracted
    print("PowerPoint Metadata Extractor")
    print("============================")
    print("This extractor can extract the following from PowerPoint presentations:")
    print("- Core properties (title, author, creation date, etc.)")
    print("- Slide dimensions and layout information") 
    print("- All shapes on each slide with:")
    print("  * Position and size (in EMU and inches)")
    print("  * Text content and formatting (fonts, colors, alignment)")
    print("  * Images with crop information")
    print("  * Tables with cell data")
    print("  * Fill, line, and shadow formatting")
    print("  * Shape-specific properties (autoshapes, groups, etc.)")
    print("- Slide notes")
    print("- Master slides and layouts")


if __name__ == "__main__":
    main()
