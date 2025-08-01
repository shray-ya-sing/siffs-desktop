"""
Create a basic PowerPoint presentation for testing purposes.
"""

from pptx import Presentation
from pptx.util import Inches

def create_basic_presentation():
    """Create a basic PowerPoint presentation with a few slides."""
    prs = Presentation()
    
    # Slide 1 - Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "For PowerPoint Table Testing"
    
    # Slide 2 - Title and Content
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sample Content Slide"
    content.text = "This slide has some sample content."
    
    # Slide 3 - Another Title and Content
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Another Sample Slide"
    content.text = "This is another test slide for reference."
    
    return prs

if __name__ == "__main__":
    # Create basic presentation
    prs = create_basic_presentation()
    
    # Save it
    filename = "basic_test_template.pptx"
    prs.save(filename)
    print(f"Created basic PowerPoint presentation: {filename}")
