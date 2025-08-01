#!/usr/bin/env python3
"""
Test script to verify the fix works with the actual PowerPoint writer.
"""

import tempfile
import os
import sys
sys.path.append('.')

def test_powerpoint_writer_with_fixed_parsing():
    """Test the PowerPoint writer with the fixed paragraph parsing."""
    print("🧪 TESTING POWERPOINT WRITER WITH FIXED PARSING")
    print("=" * 60)
    
    # Import required modules
    from ai_services.tools.read_write_functions.powerpoint.powerpoint_edit_tools import parse_markdown_powerpoint_data
    from powerpoint.editing.powerpoint_writer import PowerPointWriter
    
    # Create a temporary PowerPoint file for testing
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
        temp_path = temp_file.name
    
    print(f"Using temporary file: {temp_path}")
    
    # Create a simple PowerPoint file first
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # Create a new presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add a title
        title = slide.shapes.title
        title.text = "Test Slide"
        
        # Add a text box that we'll edit
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.name = "Text Box 2"  # Give it the name we'll reference
        text_frame = textbox.text_frame
        text_frame.text = "Original text here"
        
        # Save the presentation
        prs.save(temp_path)
        print("✅ Created test PowerPoint file")
        
    except Exception as e:
        print(f"❌ Failed to create test PowerPoint file: {e}")
        return
    
    # Test the LLM-generated markdown with problematic apostrophes
    llm_markdown = '''slide_number: slide1 | shape_name="Text Box 2", paragraphs="[{'text': 'JPMorgan Chase did not assume First Republic Bank's deposits or any other liabilities of First Republic Bank.', 'bullet_style': 'bullet', 'indent_level': 0}]"'''
    
    print(f"\\nTesting LLM markdown: {repr(llm_markdown)}")
    
    # Step 1: Parse the markdown
    try:
        parsed_data = parse_markdown_powerpoint_data(llm_markdown)
        print(f"✅ Parsed markdown successfully")
        print(f"Parsed data: {parsed_data}")
        
        if not parsed_data:
            print("❌ No data was parsed from markdown")
            return
            
    except Exception as e:
        print(f"❌ Failed to parse markdown: {e}")
        import traceback
        traceback.print_exc()
        return
    
    # Step 2: Write to PowerPoint using the writer
    try:
        writer = PowerPointWriter()
        writer.visible = True  # Show PowerPoint for inspection
        
        success, updated_shapes = writer.write_to_existing(parsed_data, temp_path)
        
        if success:
            print(f"✅ Successfully wrote to PowerPoint file")
            print(f"Updated shapes: {updated_shapes}")
            print(f"📁 PowerPoint file left open for inspection: {temp_path}")
            print(f"🔍 Please check the text box content and bullet formatting")
        else:
            print(f"❌ Failed to write to PowerPoint file")
            
    except Exception as e:
        print(f"❌ Error writing to PowerPoint: {e}")
        import traceback
        traceback.print_exc()
    finally:
        # Don't cleanup - leave PowerPoint open for inspection
        print(f"💡 PowerPoint is left open for inspection. File: {temp_path}")
        print(f"💡 The file will remain until you manually close PowerPoint or delete the temp file.")

if __name__ == "__main__":
    test_powerpoint_writer_with_fixed_parsing()
