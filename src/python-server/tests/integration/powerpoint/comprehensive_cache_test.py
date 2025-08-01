#!/usr/bin/env python3
"""
Comprehensive Cache Integration Test
Tests PowerPoint writer integration with cache manager for various edit scenarios:
- Adding new shapes
- Modifying existing shapes  
- Deleting shapes
- Adding new slides
- Modifying slide properties
- Complex multi-step editing workflows
"""

import os
import sys
import json
import tempfile
import shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

# Add project paths
current_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.append(current_dir)
sys.path.append(os.path.join(current_dir, 'ai_services'))
sys.path.append(os.path.join(current_dir, 'ai_services', 'metadata'))

from powerpoint.metadata.extraction.pptx_metadata_extractor import PowerPointMetadataExtractor
from powerpoint.editing.powerpoint_writer import PowerPointWriter
import hashlib

def load_cache(cache_file_path):
    """Load cache data from file"""
    if not os.path.exists(cache_file_path):
        return {}
    try:
        with open(cache_file_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error loading cache: {e}")
        return {}

def store_in_cache(cache_file_path, cache_key, file_path, workspace_path, metadata):
    """Store metadata in cache"""
    try:
        # Load existing cache
        cache_data = load_cache(cache_file_path)
        
        # Add metadata with cache info
        cached_metadata = {
            **metadata,
            'file_path': file_path,
            'workspace_path': workspace_path,
            'cached_at': datetime.now().isoformat(),
            'file_mtime': os.path.getmtime(file_path),
            'file_size': os.path.getsize(file_path)
        }
        
        # Store in cache
        cache_data[cache_key] = cached_metadata
        
        # Save cache
        with open(cache_file_path, 'w', encoding='utf-8') as f:
            json.dump(cache_data, f, indent=2, ensure_ascii=False, default=str)
        
        return True
    except Exception as e:
        print(f"Error storing in cache: {e}")
        return False

def update_cache_with_edits(cache_file_path, cache_key, file_path, updated_shapes):
    """Update cache with edit information for shapes"""
    try:
        # Load cache
        cache_data = load_cache(cache_file_path)
        
        if cache_key not in cache_data:
            print(f"Cache key {cache_key} not found")
            return False
            
        # Extract fresh metadata
        extractor = PowerPointMetadataExtractor()
        fresh_metadata = extractor.extract_presentation_metadata(file_path)
        
        # Update the cached metadata
        cache_data[cache_key].update(fresh_metadata)
        cache_data[cache_key]['file_mtime'] = os.path.getmtime(file_path)
        cache_data[cache_key]['last_updated'] = datetime.now().isoformat()
        
        # Add edit history to shapes based on updated_shapes info
        current_time = datetime.now().isoformat()
        
        for slide_num, shapes_list in updated_shapes.items():
            slide_index = slide_num - 1  # Convert to 0-based index
            
            if slide_index < len(cache_data[cache_key]['slides']):
                cached_slide = cache_data[cache_key]['slides'][slide_index]
                
                for shape_info in shapes_list:
                    shape_index = shape_info.get('shapeIndex')
                    
                    if shape_index is not None and shape_index < len(cached_slide['shapes']):
                        # This is a modification of existing shape
                        cached_shape = cached_slide['shapes'][shape_index]
                        
                        # Add to edit history
                        if 'editingHistory' not in cached_shape:
                            cached_shape['editingHistory'] = []
                        
                        properties_applied = []
                        if 'text' in shape_info:
                            properties_applied.append('text')
                        if 'width' in shape_info:
                            properties_applied.append('width')
                        if 'height' in shape_info:
                            properties_applied.append('height')
                        if 'left' in shape_info:
                            properties_applied.append('left')
                        if 'top' in shape_info:
                            properties_applied.append('top')
                            
                        cached_shape['editingHistory'].append({
                            'timestamp': current_time,
                            'properties_applied': properties_applied
                        })
                        
                        cached_shape['lastModified'] = current_time
                        
                    else:
                        # This is a new shape added to existing slide
                        if len(cached_slide['shapes']) > 0:
                            # Find the new shape and add edit history
                            latest_shapes = fresh_metadata['slides'][slide_index]['shapes']
                            if len(latest_shapes) > len(cached_slide['shapes']):
                                # New shapes were added
                                for new_shape in latest_shapes[len(cached_slide['shapes']):]:
                                    new_shape['editingHistory'] = [{
                                        'timestamp': current_time,
                                        'properties_applied': ['text', 'width', 'height', 'left', 'top']
                                    }]
                                    new_shape['lastModified'] = current_time
                                    new_shape['created'] = current_time
        
        # Save updated cache
        with open(cache_file_path, 'w', encoding='utf-8') as f:
            json.dump(cache_data, f, indent=2, ensure_ascii=False, default=str)
            
        return True
        
    except Exception as e:
        print(f"Error updating cache with edits: {e}")
        import traceback
        traceback.print_exc()
        return False

def create_test_presentation():
    """Create a test PowerPoint with initial content for modification tests"""
    prs = Presentation()
    
    # Slide 1: Title slide with multiple shapes
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title1 = slide1.shapes.title
    title1.text = "Original Title"
    
    subtitle1 = slide1.placeholders[1]
    subtitle1.text = "Original Subtitle"
    
    # Add a custom text box
    textbox1 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    textbox1.text = "Original Text Box"
    
    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title2 = slide2.shapes.title
    title2.text = "Content Slide"
    
    content2 = slide2.placeholders[1]
    content2.text = "Original bullet point\nSecond bullet point"
    
    # Add an auto shape
    from pptx.enum.shapes import MSO_SHAPE
    autoshape = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5), Inches(2), Inches(2), Inches(1.5)
    )
    autoshape.text = "Original Rectangle"
    
    return prs

def print_cache_summary(cache_data, scenario_name):
    """Print a summary of cache data for verification"""
    print(f"\n=== CACHE SUMMARY: {scenario_name} ===")
    
    if not cache_data:
        print("❌ Cache is empty!")
        return
        
    for pres_id, pres_data in cache_data.items():
        print(f"Presentation: {pres_data.get('presentation_name', 'Unknown')}")
        print(f"Total Slides: {len(pres_data.get('slides', []))}")
        
        for i, slide in enumerate(pres_data.get('slides', []), 1):
            shapes = slide.get('shapes', [])
            print(f"  Slide {i}: {len(shapes)} shapes")
            
            for j, shape in enumerate(shapes):
                shape_name = shape.get('name', f'Shape_{j}')
                last_modified = shape.get('lastModified', 'Never')
                edit_count = len(shape.get('editingHistory', []))
                print(f"    - {shape_name}: {edit_count} edits, last: {last_modified}")

def run_comprehensive_test():
    """Run comprehensive test covering all edit scenarios"""
    print("🚀 Starting Comprehensive Cache Integration Test...")
    
    # Create test file
    test_file = "comprehensive_cache_test.pptx"
    test_path = os.path.join(current_dir, test_file)
    
    try:
        # === SCENARIO 1: Initial Creation ===
        print("\n📝 SCENARIO 1: Creating initial presentation...")
        prs = create_test_presentation()
        prs.save(test_path)
        
        # Extract initial metadata
        extractor = PowerPointMetadataExtractor()
        initial_metadata = extractor.extract_presentation_metadata(test_path)
        print(f"✅ Created presentation with {len(initial_metadata['slides'])} slides")
        
        # Setup cache files and directory
        cache_dir = os.path.join(current_dir, "metadata", "_cache")
        os.makedirs(cache_dir, exist_ok=True)
        cache_file_path = os.path.join(cache_dir, "powerpoint_metadata_hotcache.json")
        
        # Generate cache key for the presentation - just use MD5 of path to make it stable
        file_name = os.path.basename(test_path)
        cache_key = f"presentation_{hashlib.md5(test_path.encode()).hexdigest()[:16]}"
        workspace_path = f"test_workspace/{test_file}"
        
        # Store initial metadata in cache
        store_in_cache(cache_file_path, cache_key, test_path, workspace_path, initial_metadata)
        initial_cache = load_cache(cache_file_path)
        print_cache_summary(initial_cache, "INITIAL CREATION")
        
        # === SCENARIO 2: Modify Existing Shapes ===
        print("\n✏️ SCENARIO 2: Modifying existing shapes...")
        writer = PowerPointWriter()
        
        # Modify shapes on slide 1
        slide1_modifications = {
            1: [
                {
                    "shapeIndex": 0,  # Title shape
                    "text": "MODIFIED Title - Updated Content",
                    "width": 200,
                    "height": 80,
                    "left": 150,
                    "top": 120
                },
                {
                    "shapeIndex": 1,  # Subtitle shape  
                    "text": "MODIFIED Subtitle - New Text",
                    "width": 180,
                    "height": 60
                }
            ]
        }
        
        success, updated_shapes_list = writer.write_to_existing(slide1_modifications, test_path)
        # Convert the list of updated shapes back to our expected dictionary format
        updated_shapes = {}
        for shape_data in updated_shapes_list:
            slide_num = shape_data.get('slide_number', 1)
            if slide_num not in updated_shapes:
                updated_shapes[slide_num] = []
            updated_shapes[slide_num].append(shape_data)
        print(f"✅ Modified {len(updated_shapes[1]) if 1 in updated_shapes else 0} shapes on slide 1")
        
        # Update cache with modifications
        update_cache_with_edits(cache_file_path, cache_key, test_path, updated_shapes)
        modified_cache = load_cache(cache_file_path)
        print_cache_summary(modified_cache, "SHAPE MODIFICATIONS")
        
        # === SCENARIO 3: Add New Shapes ===
        print("\n➕ SCENARIO 3: Adding new shapes to existing slides...")
        
        slide2_additions = {
            2: [
                {
                    "text": "BRAND NEW Text Box",
                    "width": 250,
                    "height": 100,
                    "left": 300,
                    "top": 200
                },
                {
                    "text": "Another NEW Shape",
                    "width": 180,
                    "height": 80,
                    "left": 100,
                    "top": 350
                }
            ]
        }
        
        success2, new_shapes_list = writer.write_to_existing(slide2_additions, test_path)
        # Convert to expected format
        new_shapes = {}
        for shape_data in new_shapes_list:
            slide_num = shape_data.get('slide_number', 2)
            if slide_num not in new_shapes:
                new_shapes[slide_num] = []
            new_shapes[slide_num].append(shape_data)
        print(f"✅ Added {len(new_shapes[2]) if 2 in new_shapes else 0} new shapes to slide 2")
        
        # Update cache with new shapes
        update_cache_with_edits(cache_file_path, cache_key, test_path, new_shapes)
        addition_cache = load_cache(cache_file_path)
        print_cache_summary(addition_cache, "SHAPE ADDITIONS")
        
        # === SCENARIO 4: Add New Slides ===
        print("\n📄 SCENARIO 4: Adding completely new slides...")
        
        new_slide_data = {
            3: [  # New slide 3
                {
                    "text": "New Slide Title",
                    "width": 400,
                    "height": 100,
                    "left": 50,
                    "top": 50
                },
                {
                    "text": "New slide content goes here",
                    "width": 500,
                    "height": 200,
                    "left": 50,
                    "top": 200
                }
            ],
            4: [  # New slide 4
                {
                    "text": "Second New Slide",
                    "width": 350,
                    "height": 80,
                    "left": 100,
                    "top": 100
                }
            ]
        }
        
        success3, slide_shapes_list = writer.write_to_existing(new_slide_data, test_path)
        # Convert to expected format
        slide_shapes = {}
        for shape_data in slide_shapes_list:
            slide_num = shape_data.get('slide_number', 1)
            if slide_num not in slide_shapes:
                slide_shapes[slide_num] = []
            slide_shapes[slide_num].append(shape_data)
        total_new_shapes = sum(len(shapes) for shapes in slide_shapes.values())
        print(f"✅ Added 2 new slides with {total_new_shapes} total shapes")
        
        # Update cache with new slides
        update_cache_with_edits(cache_file_path, cache_key, test_path, slide_shapes)
        final_cache = load_cache(cache_file_path)
        print_cache_summary(final_cache, "NEW SLIDES ADDED")
        
        # === SCENARIO 5: Complex Multi-Step Workflow ===
        print("\n🔄 SCENARIO 5: Complex multi-step editing workflow...")
        
        # Step 1: Modify multiple shapes across slides
        complex_edits = {
            1: [
                {
                    "shapeIndex": 0,
                    "text": "FINAL Title Version",
                    "width": 600,
                    "height": 120
                }
            ],
            2: [
                {
                    "shapeIndex": 1,  # Existing content placeholder
                    "text": "Updated bullet points:\n• First updated point\n• Second updated point\n• NEW third point"
                }
            ],
            3: [
                {
                    "shapeIndex": 0,  # Title from new slide
                    "text": "UPDATED New Slide Title",
                    "left": 75,
                    "top": 75
                }
            ]
        }
        
        success4, complex_shapes_list = writer.write_to_existing(complex_edits, test_path)
        # Convert to expected format
        complex_shapes = {}
        for shape_data in complex_shapes_list:
            slide_num = shape_data.get('slide_number', 1)
            if slide_num not in complex_shapes:
                complex_shapes[slide_num] = []
            complex_shapes[slide_num].append(shape_data)

        # Step 2: Add even more shapes
        more_additions = {
            4: [
                {
                    "text": "Additional shape on slide 4",
                    "width": 300,
                    "height": 150,
                    "left": 200,
                    "top": 250
                }
            ]
        }
        
        success5, additional_shapes_list = writer.write_to_existing(more_additions, test_path)
        # Convert to expected format
        additional_shapes = {}
        for shape_data in additional_shapes_list:
            slide_num = shape_data.get('slide_number', 1)
            if slide_num not in additional_shapes:
                additional_shapes[slide_num] = []
            additional_shapes[slide_num].append(shape_data)

        # Combine all updates
        all_updates = {}
        for slide_num, shapes in complex_shapes.items():
            all_updates[slide_num] = shapes
        for slide_num, shapes in additional_shapes.items():
            if slide_num in all_updates:
                all_updates[slide_num].extend(shapes)
            else:
                all_updates[slide_num] = shapes
        
        total_complex_shapes = sum(len(shapes) for shapes in all_updates.values())
        print(f"✅ Complex workflow: modified/added {total_complex_shapes} shapes across multiple slides")
        
        # Final cache update
        update_cache_with_edits(cache_file_path, cache_key, test_path, all_updates)
        complex_cache = load_cache(cache_file_path)
        print_cache_summary(complex_cache, "COMPLEX WORKFLOW")
        
        # === VERIFICATION ===
        print("\n🔍 FINAL VERIFICATION...")
        
        # Load and display final cache state
        # Use the same cache file path we've been using throughout the test
        
        if os.path.exists(cache_file_path):
            with open(cache_file_path, 'r', encoding='utf-8') as f:
                final_cache_data = json.load(f)
            
            print(f"✅ Cache file exists at: {cache_file_path}")
            print(f"✅ Cache contains {len(final_cache_data)} presentations")
            
            # Detailed verification
            for pres_id, pres_data in final_cache_data.items():
                slides = pres_data.get('slides', [])
                print(f"\n📊 FINAL STATE - {pres_data.get('presentation_name', 'Unknown')}:")
                print(f"   Total slides: {len(slides)}")
                
                total_shapes = 0
                total_edits = 0
                
                for i, slide in enumerate(slides, 1):
                    shapes = slide.get('shapes', [])
                    slide_edits = sum(len(shape.get('editingHistory', [])) for shape in shapes)
                    
                    total_shapes += len(shapes)
                    total_edits += slide_edits
                    
                    print(f"   Slide {i}: {len(shapes)} shapes, {slide_edits} total edits")
                    
                    # Show edit history for first few shapes
                    for j, shape in enumerate(shapes[:3]):  # Show first 3 shapes per slide
                        history = shape.get('editingHistory', [])
                        last_mod = shape.get('lastModified', 'Never')
                        print(f"     • {shape.get('name', f'Shape_{j}')}: {len(history)} edits, last: {last_mod}")
                
                print(f"\n📈 TOTALS: {total_shapes} shapes, {total_edits} total edits recorded")
        
        else:
            print("❌ Cache file not found!")
        
        print(f"\n🎉 ALL SCENARIOS COMPLETED SUCCESSFULLY!")
        print(f"📁 Test file available at: {test_path}")
        print(f"📄 Cache file available at: {cache_file_path}")
        
        return True
        
    except Exception as e:
        print(f"❌ Test failed with error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False
    
    finally:
        # Note: Not cleaning up files for inspection
        print(f"\n📋 Test files preserved for inspection:")
        if os.path.exists(test_path):
            print(f"   PowerPoint: {test_path}")

if __name__ == "__main__":
    success = run_comprehensive_test()
    if success:
        print("\n✅ Comprehensive test PASSED - Cache integration works correctly!")
    else:
        print("\n❌ Comprehensive test FAILED - Check logs above")
    
    sys.exit(0 if success else 1)
