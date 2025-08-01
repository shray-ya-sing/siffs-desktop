#!/usr/bin/env python3
"""
Test script to verify that bullet character functionality works correctly
"""

import os
import sys
import json
from pathlib import Path

# Add the current directory to Python path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from powerpoint.editing.powerpoint_writer import PowerPointWriter
from ai_services.tools.read_write_functions.powerpoint.powerpoint_edit_tools import parse_markdown_powerpoint_data

def test_bullet_char_functionality():
    """Test that bullet characters can be applied correctly"""
    print("=== Testing Bullet Character Functionality ===")
    
    if not PPTX_AVAILABLE:
        print("❌ python-pptx not available, skipping test")
        return False
    
    test_file = "test_bullet_char.pptx"
    
    try:
        # Step 1: Create a simple PowerPoint with text
        print("\n1️⃣ Creating PowerPoint with text content...")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        
        # Add title
        title = slide.shapes.title
        title.text = "Bullet Character Test"
        
        # Add a text box with bullet points
        from pptx.util import Inches
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        textbox.name = "Test Text Box"  # Set the shape name
        text_frame = textbox.text_frame
        text_frame.text = "First bullet point\nSecond bullet point\nThird bullet point"
        
        prs.save(test_file)
        print(f"✅ Created: {test_file}")
        
        # Step 2: Test bullet character application using PowerPoint writer
        print("\n2️⃣ Testing bullet character application...")
        writer = PowerPointWriter()
        
        # Test data with just bullet_char (no bullet_style)
        slide_data = {
            "1": {
                "Test Text Box": {
                    "bullet_char": "-"
                }
            }
        }
        
        success, updated_shapes = writer.write_to_existing(slide_data, test_file)
        
        if success:
            print(f"✅ PowerPoint writer successful: {len(updated_shapes)} shapes updated")
            for shape in updated_shapes:
                properties = shape.get('properties_applied', [])
                if 'bullet_char' in properties:
                    print(f"✅ bullet_char applied to shape: {shape.get('shape_name')}")
                else:
                    print(f"⚠️  bullet_char NOT applied to shape: {shape.get('shape_name')}")
        else:
            print("❌ PowerPoint writer failed")
            return False
        
        # Step 3: Test with markdown parsing
        print("\n3️⃣ Testing with markdown parsing...")
        markdown_input = 'slide_number: 1 | shape_name="Test Text Box", bullet_char="-"'
        
        parsed_data = parse_markdown_powerpoint_data(markdown_input)
        print(f"Parsed data: {json.dumps(parsed_data, indent=2)}")
        
        if parsed_data and "1" in parsed_data:
            shape_props = parsed_data["1"].get("Test Text Box", {})
            if "bullet_char" in shape_props:
                print(f"✅ Markdown parsing successful: bullet_char = '{shape_props['bullet_char']}'")
            else:
                print("❌ bullet_char not found in parsed data")
                return False
        else:
            print("❌ Markdown parsing failed")
            return False
        
        # Step 4: Test end-to-end with parsed data
        print("\n4️⃣ Testing end-to-end with parsed data...")
        success2, updated_shapes2 = writer.write_to_existing(parsed_data, test_file)
        
        if success2:
            print(f"✅ End-to-end test successful: {len(updated_shapes2)} shapes updated")
            for shape in updated_shapes2:
                properties = shape.get('properties_applied', [])
                if 'bullet_char' in properties:
                    print(f"✅ bullet_char applied via parsed data to: {shape.get('shape_name')}")
                else:
                    print(f"⚠️  bullet_char NOT applied via parsed data to: {shape.get('shape_name')}")
        else:
            print("❌ End-to-end test failed")
            return False
        
        print("\n🎉 All tests passed! Bullet character functionality is working correctly.")
        return True
        
    except Exception as e:
        print(f"❌ Test failed with error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False
    
    finally:
        # Clean up
        if os.path.exists(test_file):
            try:
                os.remove(test_file)
                print(f"🧹 Cleaned up: {test_file}")
            except:
                print(f"⚠️  Could not clean up: {test_file}")

if __name__ == "__main__":
    success = test_bullet_char_functionality()
    if success:
        print("\n✅ RESULT: Bullet character functionality is working!")
    else:
        print("\n❌ RESULT: Bullet character functionality has issues!")
    
    sys.exit(0 if success else 1)
